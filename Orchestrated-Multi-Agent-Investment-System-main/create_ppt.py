from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # Define Colors
    NEON_GREEN = RGBColor(57, 255, 20)
    CYAN = RGBColor(0, 255, 255)
    WHITE = RGBColor(255, 255, 255)
    DARK_BLUE = RGBColor(10, 10, 40)
    
    # Helper to apply background
    def apply_background(slide, image_path):
        if os.path.exists(image_path):
            left = top = 0
            pic = slide.shapes.add_picture(image_path, left, top, width=prs.slide_width, height=prs.slide_height)
            # Move to back
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)

    # Helper to set text formatting
    def set_text_format(font_obj, font_size=18, bold=False, color=WHITE, font_name='Arial'):
        font_obj.size = Pt(font_size)
        font_obj.name = font_name
        font_obj.bold = bold
        if color:
            font_obj.color.rgb = color

    # 1. Title Slide
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, 'ppt_background.png')
    
    # Add Poster Image
    if os.path.exists('project_poster_v2.png'):
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(5.5)
        slide.shapes.add_picture('project_poster_v2.png', left, top, height=height)
    
    # Title Text (Right side)
    txBox = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4.5), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "ORCHESTRATED\nMULTI-AGENT\nINVESTMENT SYSTEM"
    p.alignment = PP_ALIGN.LEFT
    set_text_format(p.font, 44, True, NEON_GREEN, 'Impact')
    
    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(5), Inches(4.5), Inches(4.5), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "AI-Powered Platform for Automated, Comprehensive Financial Analysis"
    set_text_format(p.font, 20, False, CYAN)

    # Team
    txBox = slide.shapes.add_textbox(Inches(5), Inches(6), Inches(4.5), Inches(1.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Team Members:"
    set_text_format(p.font, 18, True, WHITE)
    
    members = [
        "Adarsh Prabhu (1DT22AI004)",
        "Atul Mayank (1DT22AI007)",
        "Ayush Poddar (1DT22AI008)",
        "Dinesh Krishan (1DT22AI020)"
    ]
    for m in members:
        p = tf.add_paragraph()
        p.text = m
        set_text_format(p.font, 16, False, WHITE)

    # Common Slide Function
    def add_content_slide(title_text, content_lines):
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        apply_background(slide, 'ppt_background.png')
        
        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = title_text
        set_text_format(p.font, 40, True, NEON_GREEN, 'Impact')
        
        # Content
        top = Inches(1.8)
        for line in content_lines:
            txBox = slide.shapes.add_textbox(Inches(1), top, Inches(8.5), Inches(0.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = "• " + line
            set_text_format(p.font, 24, False, WHITE)
            top += Inches(0.8)

    # 2. Introduction
    add_content_slide("Introduction", [
        "Manual investment analysis is time-consuming and prone to bias.",
        "Difficulty in aggregating data from multiple sources (Technical, Fundamental, News).",
        "Need for a unified, automated system for comprehensive financial insights."
    ])

    # 3. System Overview
    add_content_slide("System Overview", [
        "An advanced AI platform utilizing multiple specialized agents.",
        "Automates the entire investment analysis pipeline.",
        "Provides data-driven, holistic investment recommendations."
    ])

    # 4. Architecture: The Agents
    add_content_slide("System Architecture", [
        "Monitor Agent: Continuously tracks market data and stock prices.",
        "Planner Agent: Creates initial investment strategies based on user goals.",
        "Predictor Agent: Validates plans using predictive models and risk analysis.",
        "Aggregator: Synthesizes outputs to generate the final report."
    ])

    # 5. Key Features
    add_content_slide("Key Features", [
        "Interactive Web Interface (FastAPI + HTML/CSS/JS)",
        "Real-time Market Data Integration",
        "Comprehensive Risk Assessment",
        "Personalized Investment Portfolios",
        "Automated Validation and Reporting"
    ])

    # 6. Tech Stack
    add_content_slide("Technology Stack", [
        "Backend: Python, FastAPI",
        "AI/ML: Large Language Models (LLMs), LangChain",
        "Frontend: HTML5, CSS3, JavaScript",
        "Data: Financial APIs (e.g., yfinance)",
        "Visualization: Matplotlib"
    ])

    # 7. Workflow
    add_content_slide("System Workflow", [
        "1. User Input: Customer profile and investment goals.",
        "2. Data Retrieval: Monitor Agent fetches market data.",
        "3. Strategy Formation: Planner Agent proposes a portfolio.",
        "4. Validation: Predictor Agent assesses risk and returns.",
        "5. Final Output: Comprehensive investment plan presented to user."
    ])

    # 8. Future Scope
    add_content_slide("Future Scope", [
        "Mobile Application Development",
        "Integration with live brokerage accounts for auto-execution",
        "Advanced Sentiment Analysis from social media news",
        "Multi-currency support for global markets"
    ])

    # 9. Conclusion
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, 'ppt_background.png')
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Thank You!"
    p.alignment = PP_ALIGN.CENTER
    set_text_format(p.font, 60, True, NEON_GREEN, 'Impact')
    
    p = tf.add_paragraph()
    p.text = "Any Questions?"
    p.alignment = PP_ALIGN.CENTER
    set_text_format(p.font, 32, False, CYAN)

    prs.save('Project_Presentation_v2.pptx')
    print("Presentation saved as 'Project_Presentation_v2.pptx'")

if __name__ == "__main__":
    create_presentation()
